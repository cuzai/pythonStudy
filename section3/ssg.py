import requests
from bs4 import BeautifulSoup
import os, errno
from urllib import request
from pptx import Presentation

SSG = "http://www.ssg.com/event/eventMain.ssg?gnb=event"

def downText(eventList) :
    path = "C:/Users/cuzai/Desktop/Web_Crawling/SSG"
    txtPath = path + "/SSG.txt"
    with open(txtPath, 'w') as f:
      for n, i in enumerate(eventList, 1) :
            try :
              os.makedirs(path)
            except OSError as e :
                 if e.errno != errno.EEXIST :
                    raise
            f.write(str(n) + "." + i.select_one('.eo_tit > strong').string + "\n\n")
    print("downText finished")

def downImg(eventList) :
    path = "C:/Users/cuzai/Desktop/Web_Crawling/SSG"
    for n, i in enumerate(eventList, 1) :
        try :
          os.makedirs(path)
        except OSError as e :
             if e.errno != errno.EEXIST :
                raise
        try :
            title = str(n)+ ". " + i.select_one('.eo_tit > strong').string
            imgPath = os.path.join(path,  title + ".png")
            request.urlretrieve(i.select_one('.thmb > img')['src'], imgPath)
        except OSError:
            if ':' or '*'or '?'or '"'or '<'or '>'or '|' in title :
                title.replace(":", "-").replace("*", "_").replace("?", "..").replace('"', "'").replace("<", "(").replace(">", ")")
                #imgPath = os.path.join(path, title + ".png")
                #request.urlretrieve(i.select_one('.thmb > img')['src'], imgPath)
    print("downImg finished")

def toPptx(path) :
    ppt = Presentation(path)
    slide = ppt.slides[0]   #첫번째 슬라이드를 사용
    print(slide.shapes[0].has_table)



with requests.Session() as s :
    url = s.get(SSG)
    #print(url.text)
    soup = BeautifulSoup(url.content, 'html.parser')
    #print(soup)
    eventList = soup.select('#link1 li')
    #print(eventList)
    downText(eventList)
    downImg(eventList)
    #toPptx('C:/Users/cuzai/Desktop/Web_Crawling/SSG/promo.pptx')



ppt = Presentation('C:/Users/cuzai/Desktop/Web_Crawling/SSG/promo.pptx')
slide = ppt.slides[0]
print(slide.shapes[0].has_table)
table = slide.shapes[0].table
cell11 = table.cell(1, 1).text_frame

cell11.clear()
p = cell11.paragraphs[0]
#run = p.add_run()
#run2 = p.add_run()
#run.text = '1\n'
#run2.text = '2'

myList = ['가', '나', '다', '라', '마']
for i in myList :
    run = p.add_run()
    run.text = i

#ppt.save('C:/Users/cuzai/Desktop/Web_Crawling/SSG/promo.pptx')